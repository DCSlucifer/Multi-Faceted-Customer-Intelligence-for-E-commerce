"""Generate the 20-slide presentation deck for Team 4's final project.

Produces `reports/slides.pptx` populated with real numbers and figures from
the executed notebooks. Run after notebooks 01..06 have completed.

Usage:
    python -m src.build_slides
"""
from __future__ import annotations

from pathlib import Path
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from .submission_snapshot import load_submission_snapshot
from .utils import REPORTS_DIR, FIGURES_DIR, get_logger

logger = get_logger(__name__)

# --------------------------------------------------------------------------- #
# Style constants                                                              #
# --------------------------------------------------------------------------- #
PRIMARY   = RGBColor(0x26, 0x46, 0x53)   # dark teal
ACCENT    = RGBColor(0xE7, 0x6F, 0x51)   # warm orange
INK       = RGBColor(0x1F, 0x1F, 0x1F)
SUBINK    = RGBColor(0x55, 0x55, 0x55)
BG_SOFT   = RGBColor(0xF6, 0xF1, 0xE7)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# --------------------------------------------------------------------------- #
# Helper builders                                                              #
# --------------------------------------------------------------------------- #
def add_title_bar(slide, title: str, subtitle: str | None = None):
    """Top title bar, present on every content slide."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.7))
    bar.fill.solid(); bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    tf = bar.text_frame; tf.margin_left = Inches(0.4); tf.margin_top = Inches(0.12)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(24); r.font.bold = True; r.font.color.rgb = RGBColor(0xFF,0xFF,0xFF)
    if subtitle:
        tx = slide.shapes.add_textbox(Inches(0.4), Inches(0.78), Inches(12), Inches(0.4))
        tf = tx.text_frame; p = tf.paragraphs[0]
        r = p.add_run(); r.text = subtitle
        r.font.size = Pt(14); r.font.color.rgb = SUBINK; r.font.italic = True


def add_textbox(slide, left, top, width, height, text, size=14,
                bold=False, color=INK, align=PP_ALIGN.LEFT):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return tx


def add_bullets(slide, left, top, width, height, items, size=16, bold_first_word=False):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = "•  " + item
        r.font.size = Pt(size); r.font.color.rgb = INK
        p.space_after = Pt(8)


def add_image(slide, path: Path, left, top, width=None, height=None):
    if not path.exists():
        # Fallback: a labelled placeholder rectangle so the deck still renders.
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top,
                                       width or Inches(6), height or Inches(4))
        rect.fill.solid(); rect.fill.fore_color.rgb = BG_SOFT
        tf = rect.text_frame
        tf.text = f"[missing figure]\n{path.name}"
        return
    if width and height:
        slide.shapes.add_picture(str(path), left, top, width=width, height=height)
    elif width:
        slide.shapes.add_picture(str(path), left, top, width=width)
    elif height:
        slide.shapes.add_picture(str(path), left, top, height=height)
    else:
        slide.shapes.add_picture(str(path), left, top)


def add_footer(slide, page_no: int, total: int = 20):
    add_textbox(slide, Inches(0.4), Inches(7.05), Inches(6), Inches(0.35),
                "Team 4 — Multi-Faceted Customer Intelligence (Online Retail II)",
                size=10, color=SUBINK)
    add_textbox(slide, Inches(12.5), Inches(7.05), Inches(0.7), Inches(0.35),
                f"{page_no} / {total}", size=10, color=SUBINK, align=PP_ALIGN.RIGHT)


def add_speaker_notes(slide, text: str):
    slide.notes_slide.notes_text_frame.text = text


# --------------------------------------------------------------------------- #
# Slide content                                                                #
# --------------------------------------------------------------------------- #
def build_deck() -> Presentation:
    snapshot = load_submission_snapshot()
    tokens = snapshot.tokens()

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]   # totally blank

    # ---------- 1. Title -------------------------------------------------- #
    s = prs.slides.add_slide(blank)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.0), SLIDE_W, Inches(3.5))
    bar.fill.solid(); bar.fill.fore_color.rgb = PRIMARY; bar.line.fill.background()
    add_textbox(s, Inches(0.5), Inches(2.2), Inches(12.3), Inches(1.0),
                "Multi-Faceted Customer Intelligence", size=42, bold=True,
                color=RGBColor(0xFF,0xFF,0xFF), align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.5), Inches(3.2), Inches(12.3), Inches(0.8),
                "An End-to-End Data Mining Pipeline on Online Retail II",
                size=22, color=RGBColor(0xFF,0xFF,0xFF), align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.5),
                "Data Mining 2026 — Final Project", size=16,
                color=RGBColor(0xFF,0xFF,0xFF), align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.5), Inches(5.7), Inches(12.3), Inches(0.4),
                "Team 4 — University of Information Technology, VNU-HCM",
                size=18, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.4),
                "Nguyen Minh Cuong (22520177)   •   Vo Thanh Danh (22520201)   "
                "•   Nguyen Vinh Dat (22520228)   •   Nguyen Huu Dinh (22520251)",
                size=13, color=SUBINK, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
                "02 June 2026", size=12, color=SUBINK, align=PP_ALIGN.CENTER)

    # ---------- 2. Problem & motivation ----------------------------------- #
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "The Problem", "Three siloed questions every retailer asks")
    add_bullets(s, Inches(0.6), Inches(1.4), Inches(7.5), Inches(4),
                ["WHO are my customers? — segmentation",
                 "WHO will leave? — churn prediction",
                 "WHAT do they buy together? — basket analysis"], size=20)
    add_textbox(s, Inches(0.6), Inches(4.3), Inches(8), Inches(0.5),
                "Solved in isolation, these analyses miss the joint signal that drives retention ROI.",
                size=15, color=SUBINK)
    # Big stat card on the right
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(8.7), Inches(1.4), Inches(4.2), Inches(4.5))
    card.fill.solid(); card.fill.fore_color.rgb = BG_SOFT; card.line.fill.background()
    add_textbox(s, Inches(8.7), Inches(1.6), Inches(4.2), Inches(0.6),
                "Why this matters", size=16, bold=True, color=PRIMARY,
                align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(8.7), Inches(2.4), Inches(4.2), Inches(2.0),
                "+5%", size=72, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(8.7), Inches(4.1), Inches(4.2), Inches(0.5),
                "customer retention", size=14, color=INK, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(8.7), Inches(4.6), Inches(4.2), Inches(1.5),
                "drives 25–95% profit growth\n(Bain & Co.)", size=13,
                color=SUBINK, align=PP_ALIGN.CENTER)
    add_footer(s, 2)
    add_speaker_notes(s, "Open with the three questions. Emphasise the cost asymmetry: "
                         "retaining is 5-25x cheaper than acquiring.")

    # ---------- 3. Dataset ------------------------------------------------ #
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Dataset — UCI Online Retail II",
                  "UK-centric online retailer, 2009-12 to 2011-12")
    # Stat row
    stats = [("1.07M", "raw transactions"), ("805K", "after cleaning"),
             ("5,878", "unique customers"), ("4,631", "stock codes"),
             ("41", "countries"), ("2 yr", "time span")]
    for i, (val, lab) in enumerate(stats):
        col = i % 3; row = i // 3
        l = Inches(0.6 + col * 4.2); t = Inches(1.5 + row * 1.6)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, Inches(4.0), Inches(1.4))
        card.fill.solid(); card.fill.fore_color.rgb = BG_SOFT; card.line.fill.background()
        add_textbox(s, l, t + Inches(0.15), Inches(4.0), Inches(0.7),
                    val, size=32, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
        add_textbox(s, l, t + Inches(0.85), Inches(4.0), Inches(0.4),
                    lab, size=13, color=SUBINK, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.6), Inches(5.0), Inches(12), Inches(0.4),
                "Columns (8, all transaction-level): Invoice · StockCode · Description · Quantity · Price · InvoiceDate · CustomerID · Country",
                size=13, color=SUBINK)
    add_textbox(s, Inches(0.6), Inches(5.45), Inches(12), Inches(0.45),
                "No customer attributes (name / contact / demographics) exist beyond the CustomerID key "
                "— a null CustomerID therefore carries no recoverable customer information.",
                size=12, bold=True, color=PRIMARY)
    add_textbox(s, Inches(0.6), Inches(5.95), Inches(12), Inches(0.4),
                "Source: https://archive.ics.uci.edu/dataset/502/online+retail+ii",
                size=12, color=SUBINK)
    add_footer(s, 3)

    # ---------- 4. System architecture ----------------------------------- #
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "System Architecture",
                  "Users & automation → entry points → pipeline services → artifact stores")
    add_image(s, FIGURES_DIR / "00_architecture.png", Inches(0.5), Inches(1.7), width=Inches(12.3))
    add_footer(s, 4)
    add_speaker_notes(s, "One coherent, reproducible pipeline: all four course topics + a DL "
                         "extension, driven by a CLI/submission-builder, guarded by tests + a "
                         "validator, with the report/slides/README generated from one metrics snapshot.")

    # ---------- 5. Preprocessing ----------------------------------------- #
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Preprocessing",
                  "From raw 1.07M rows to clean 805K rows — only 24.5% dropped")
    add_bullets(s, Inches(0.5), Inches(1.25), Inches(6.0), Inches(2.8),
                ["Drop 22.8% rows with null CustomerID",
                 "Remove cancellation invoices (prefix 'C')",
                 "Drop non-positive Quantity / Price",
                 "Winsorize Quantity & Price at 1%/99% quantiles",
                 "Derive TotalAmount = Quantity × Price",
                 "Compute RFM + behavioral features per customer",
                 "Label churn: no activity in 90 d after 2011-09-01"], size=13)
    # Cleaning breakdown table — deterministic dataset facts (report Table II).
    # Filters overlap, so per-filter shares need not sum to the total.
    clean_rows = [["Filter removed", "% of raw rows"],
                  ["Missing CustomerID", "22.8%"],
                  ["Cancellation invoices (C)", "1.8%"],
                  ["Non-positive Qty / Price", "2.4%"],
                  ["Total rows removed", "24.5%"]]
    tbl = s.shapes.add_table(rows=len(clean_rows), cols=2,
                             left=Inches(6.8), top=Inches(1.35),
                             width=Inches(6.1), height=Inches(2.2)).table
    for i, row in enumerate(clean_rows):
        for j, v in enumerate(row):
            cell = tbl.cell(i, j); cell.text = v
            p = cell.text_frame.paragraphs[0]
            for r_ in p.runs:
                r_.font.size = Pt(13)
                r_.font.bold = (i == 0) or (i == len(clean_rows) - 1)
                r_.font.color.rgb = RGBColor(0xFF,0xFF,0xFF) if i == 0 else INK
            if i == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = PRIMARY
    add_textbox(s, Inches(6.8), Inches(3.65), Inches(6.1), Inches(0.9),
                "Dropping null-CustomerID rows is mandatory (RFM / churn are "
                "per-customer) yet costs only 3.5% of units & 13.7% of revenue "
                "— almost no purchasing signal lost.",
                size=12, bold=True, color=PRIMARY)
    # Distributions figure is 3.2:1 — placed below, sized to clear the table/note.
    add_image(s, FIGURES_DIR / "01_distributions.png",
              Inches(3.37), Inches(4.6), width=Inches(6.6))
    add_textbox(s, Inches(0.5), Inches(6.78), Inches(12.3), Inches(0.25),
                "Quantity / Price / TotalAmount distributions after winsorization",
                size=11, color=SUBINK, align=PP_ALIGN.CENTER)
    add_footer(s, 5)

    # ---------- 6. EDA: revenue + Pareto --------------------------------- #
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "EDA — Revenue & Customer Pareto",
                  "Q4 seasonality + top 20% of customers drive most revenue")
    add_image(s, FIGURES_DIR / "01_monthly_revenue.png",
              Inches(0.4), Inches(1.4), width=Inches(6.3))
    add_image(s, FIGURES_DIR / "01_pareto.png",
              Inches(6.8), Inches(1.4), width=Inches(6.3))
    add_footer(s, 6)

    # ---------- 7. EDA: products / countries ----------------------------- #
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "EDA — Top Products & Export Countries",
                  "Long-tail product mix; concentrated international demand")
    add_image(s, FIGURES_DIR / "01_top_products_countries.png",
              Inches(0.6), Inches(1.3), width=Inches(12.0))
    add_footer(s, 7)

    # ---------- 8. Clustering — methods ---------------------------------- #
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Customer Segmentation — Three Algorithms Compared",
                  "K-means vs DBSCAN vs AGNES on log-transformed RFM")
    val = pd.read_csv(REPORTS_DIR / "clustering_validation.csv")
    # Build a small table
    rows = [["Algorithm", "k", "Silhouette ↑", "Davies-B ↓", "Calinski-H ↑"]]
    for _, r in val.iterrows():
        rows.append([r["Algorithm"], str(int(r["n_clusters"])),
                     f"{r['Silhouette']:.3f}", f"{r['Davies-Bouldin']:.3f}",
                     f"{r['Calinski-Harabasz']:.0f}"])
    tbl = s.shapes.add_table(rows=len(rows), cols=5,
                             left=Inches(0.6), top=Inches(1.4),
                             width=Inches(7.5), height=Inches(2.6)).table
    for i, row in enumerate(rows):
        for j, val_s in enumerate(row):
            cell = tbl.cell(i, j); cell.text = val_s
            p = cell.text_frame.paragraphs[0]
            for r_ in p.runs:
                r_.font.size = Pt(13)
                r_.font.bold = (i == 0)
                r_.font.color.rgb = RGBColor(0xFF,0xFF,0xFF) if i == 0 else INK
            if i == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = PRIMARY
    add_bullets(s, Inches(0.6), Inches(4.4), Inches(7.5), Inches(3),
                ["K-means k=2 wins silhouette (0.42) — best partitional fit",
                 "DBSCAN fragments into 13 + noise — wrong fit for heavy-tail",
                 "AGNES Ward close second; chosen k=2 confirmed by dendrogram"],
                size=14)
    add_image(s, FIGURES_DIR / "02_kmeans_sweep.png",
              Inches(8.4), Inches(1.5), width=Inches(4.7))
    add_footer(s, 8)
    add_speaker_notes(s, "Emphasise that DBSCAN's negative silhouette is the textbook signal that "
                         "density-based clustering fails on heavy-tail RFM data.")

    # ---------- 9. Clustering — profile ---------------------------------- #
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Two Customer Segments — Crystal Clear Interpretation",
                  "Dormant (54%) vs Active loyal (46%)")
    add_image(s, FIGURES_DIR / "02_cluster_profile.png",
              Inches(0.5), Inches(1.4), width=Inches(7.5))
    # Segment cards
    for i, (name, recency, freq, money, color) in enumerate([
        ("Dormant",      "317 d", "1.85", "£496",   ACCENT),
        ("Active loyal", "63 d",  "11.55", "£5,206", PRIMARY),
    ]):
        l = Inches(8.4); t = Inches(1.4 + i * 2.5)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, Inches(4.5), Inches(2.3))
        card.fill.solid(); card.fill.fore_color.rgb = color; card.line.fill.background()
        white = RGBColor(0xFF,0xFF,0xFF)
        add_textbox(s, l, t + Inches(0.15), Inches(4.5), Inches(0.5),
                    name, size=18, bold=True, color=white, align=PP_ALIGN.CENTER)
        add_textbox(s, l, t + Inches(0.7), Inches(4.5), Inches(1.6),
                    f"Recency  {recency}\nFrequency  {freq}\nMonetary  {money}",
                    size=15, color=white, align=PP_ALIGN.CENTER)
    add_footer(s, 9)

    # ---------- 10. Clustering — viz ------------------------------------ #
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "2-D Projection — PCA & t-SNE",
                  "Clean separation along the first principal component")
    add_image(s, FIGURES_DIR / "02_pca_tsne.png",
              Inches(0.6), Inches(1.4), width=Inches(12.0))
    add_footer(s, 10)

    # ---------- 11. Classification — leaderboard ------------------------ #
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Churn Classification — 7 Models Compared",
                  "5-fold CV + held-out test; SMOTE inside folds")
    cls = pd.read_csv(REPORTS_DIR / "classical_results.csv").sort_values("AUC", ascending=False)
    mlp = pd.read_csv(REPORTS_DIR / "mlp_vs_classical.csv").iloc[1]
    best = cls.iloc[0]
    runner = cls.iloc[1]
    auc_gap = best["AUC"] - runner["AUC"]
    rows = [["Model", "AUC", "F1", "Precision", "Recall"]]
    for _, r in cls.iterrows():
        rows.append([r["Model"], f"{r['AUC']:.3f}", f"{r['F1']:.3f}",
                     f"{r['Precision']:.3f}", f"{r['Recall']:.3f}"])
    rows.append(["MLP (PyTorch)", f"{mlp['AUC']:.3f}", f"{mlp['F1']:.3f}",
                 f"{mlp['Precision']:.3f}", f"{mlp['Recall']:.3f}"])
    tbl = s.shapes.add_table(rows=len(rows), cols=5,
                             left=Inches(0.6), top=Inches(1.4),
                             width=Inches(7.5), height=Inches(4.3)).table
    for i, row in enumerate(rows):
        for j, v in enumerate(row):
            cell = tbl.cell(i, j); cell.text = v
            p = cell.text_frame.paragraphs[0]
            for r_ in p.runs:
                r_.font.size = Pt(12)
                r_.font.bold = (i == 0)
                r_.font.color.rgb = RGBColor(0xFF,0xFF,0xFF) if i == 0 else INK
            if i == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = PRIMARY
    add_image(s, FIGURES_DIR / "03_roc_classical.png",
              Inches(8.4), Inches(1.5), width=Inches(4.7))
    add_footer(s, 11)

    # ---------- 12. Classification — best + confusion --------------------- #
    s = prs.slides.add_slide(blank)
    add_title_bar(s, f"Best Model — {best['Model']} (AUC {best['AUC']:.3f})",
                  f"{runner['Model']} within {auc_gap:.3f} · MLP AUC {mlp['AUC']:.3f}")
    add_image(s, FIGURES_DIR / "03_cm_best.png",
              Inches(0.7), Inches(1.4), width=Inches(5.5))
    add_bullets(s, Inches(6.5), Inches(1.5), Inches(6.5), Inches(5),
                [f"AUC {best['AUC']:.4f} · F1 {best['F1']:.4f} · Recall {best['Recall']:.4f}",
                 f"MLP test: AUC {mlp['AUC']:.4f} · Precision {mlp['Precision']:.4f} · Recall {mlp['Recall']:.4f}",
                 "Classical and MLP converge — feature ceiling on RFM+behavioural",
                 "→ classical models preferred (cheaper, interpretable)",
                 "→ DL would need richer signals (sequences, sessions) to pull ahead"],
                size=15)
    add_footer(s, 12)

    # ---------- 13. SHAP -------------------------------------------------- #
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Explainability — SHAP on RandomForest",
                  "Recency dominates; behavioural features add the rest")
    add_image(s, FIGURES_DIR / "05_shap_beeswarm.png",
              Inches(0.5), Inches(1.4), width=Inches(6.5))
    add_image(s, FIGURES_DIR / "03_rf_importance.png",
              Inches(7.3), Inches(1.4), width=Inches(5.7))
    add_footer(s, 13)
    add_speaker_notes(s, "Mention the Recency-churn structural correlation here — "
                         "we revisit it in limitations.")

    # ---------- 14. Association rules ------------------------------------ #
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Market Basket — Apriori vs FP-Growth",
                  "Same 242 itemsets, 17 rules at lift ≥ 1.5 — FP-Growth 1.5× faster")
    rt = pd.read_csv(REPORTS_DIR / "apriori_vs_fpgrowth.csv")
    rows = [["Algorithm", "Itemsets", "Rules", "Runtime (s)"]]
    for _, r in rt.iterrows():
        rows.append([r["Algorithm"], str(r["Freq itemsets"]),
                     str(r["Rules (lift>=1.5)"]), f"{r['Runtime (s)']:.2f}"])
    tbl = s.shapes.add_table(rows=len(rows), cols=4,
                             left=Inches(0.5), top=Inches(1.4),
                             width=Inches(6.0), height=Inches(1.8)).table
    for i, row in enumerate(rows):
        for j, v in enumerate(row):
            cell = tbl.cell(i, j); cell.text = v
            p = cell.text_frame.paragraphs[0]
            for r_ in p.runs:
                r_.font.size = Pt(13); r_.font.bold = (i == 0)
                r_.font.color.rgb = RGBColor(0xFF,0xFF,0xFF) if i == 0 else INK
            if i == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = PRIMARY
    # Bar chart is 1.5:1 — constrain by HEIGHT so it doesn't overflow into the footer
    add_image(s, FIGURES_DIR / "04_apriori_vs_fpgrowth.png",
              Inches(1.0), Inches(3.5), height=Inches(3.3))
    add_image(s, FIGURES_DIR / "04_rule_scatter.png",
              Inches(7.0), Inches(1.4), width=Inches(6.0))
    add_footer(s, 14)

    # ---------- 15. Top rules + network --------------------------------- #
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Top Association Rules",
                  "Strongest signals = collectible sets + color-variant pairs")
    rules = pd.read_csv(REPORTS_DIR / "association_rules.csv").head(5)
    rows = [["Antecedent", "Consequent", "Conf.", "Lift"]]
    for _, r in rules.iterrows():
        rows.append([r["antecedents_desc"][:34], r["consequents_desc"][:34],
                     f"{r['confidence']:.2f}", f"{r['lift']:.1f}"])
    tbl = s.shapes.add_table(rows=len(rows), cols=4,
                             left=Inches(0.5), top=Inches(1.5),
                             width=Inches(7.5), height=Inches(2.8)).table
    for i, row in enumerate(rows):
        for j, v in enumerate(row):
            cell = tbl.cell(i, j); cell.text = v
            p = cell.text_frame.paragraphs[0]
            for r_ in p.runs:
                r_.font.size = Pt(11); r_.font.bold = (i == 0)
                r_.font.color.rgb = RGBColor(0xFF,0xFF,0xFF) if i == 0 else INK
            if i == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = PRIMARY
    add_image(s, FIGURES_DIR / "04_rule_network.png",
              Inches(8.3), Inches(1.4), width=Inches(4.8))
    add_textbox(s, Inches(0.5), Inches(4.6), Inches(7.5), Inches(2.1),
                "Pattern: customers who buy one item from a colour or "
                "themed set frequently complete the set. Operational use: "
                "post-purchase recommendations and bundled retention offers.\n"
                "Rules are mined on customer-attributed invoices only — consistent "
                "with the per-segment rules, and excluding anonymous wholesale/bulk baskets.",
                size=13, color=INK)
    add_footer(s, 15)

    # ---------- 16. DL extension ---------------------------------------- #
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Deep Learning Extension",
                  "MLP for churn + Autoencoder for anomaly detection")
    add_image(s, FIGURES_DIR / "05_mlp_curves.png",
              Inches(0.5), Inches(1.4), width=Inches(6.3))
    add_image(s, FIGURES_DIR / "05_autoencoder.png",
              Inches(6.9), Inches(1.4), width=Inches(6.3))
    add_textbox(s, Inches(0.5), Inches(5.0), Inches(12.5), Inches(1.5),
                f"MLP [64,32,16] with SMOTE + early stopping reaches test AUC {mlp['AUC']:.3f} "
                f"(F1 {mlp['F1']:.3f}, precision {mlp['Precision']:.3f}) — on par with the best "
                f"classical model. Autoencoder (bottleneck=4) flags the top 5% of customers by "
                "reconstruction error as anomalies — candidates for VIP review or fraud screening.",
                size=14, color=INK)
    add_footer(s, 16)

    # ---------- 17. HEADLINE — churn per segment ------------------------ #
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "★ HEADLINE FINDING — Churn Rate by Segment",
                  "67-percentage-point gap between segments")
    add_image(s, FIGURES_DIR / "06_churn_by_cluster.png",
              Inches(0.5), Inches(1.4), height=Inches(3.4))
    # Segment churn + value table (report Table VI), read from artifacts.
    ccr = pd.read_csv(REPORTS_DIR / "cluster_churn_rate.csv").sort_values(
        "avg_recency", ascending=False)
    seg_names = ["Dormant", "Active loyal"]
    seg_rows = [["Segment", "n", "Churn %", "Avg Recency", "Avg Monetary (£)"]]
    for name, (_, r) in zip(seg_names, ccr.iterrows()):
        seg_rows.append([name, f"{int(r['n_customers']):,}", f"{r['churn_rate']:.1f}",
                         f"{r['avg_recency']:.1f}", f"{r['avg_monetary']:,.1f}"])
    churn_overall = (ccr['n_customers'] * ccr['churn_rate']).sum() / ccr['n_customers'].sum()
    seg_rows.append(["Overall", f"{int(ccr['n_customers'].sum()):,}",
                     f"{churn_overall:.1f}", "—", "—"])
    tbl = s.shapes.add_table(rows=len(seg_rows), cols=5,
                             left=Inches(0.5), top=Inches(5.0),
                             width=Inches(7.8), height=Inches(1.6)).table
    for i, row in enumerate(seg_rows):
        for j, v in enumerate(row):
            cell = tbl.cell(i, j); cell.text = v
            p = cell.text_frame.paragraphs[0]
            for r_ in p.runs:
                r_.font.size = Pt(12)
                r_.font.bold = (i == 0) or (i == len(seg_rows) - 1)
                r_.font.color.rgb = RGBColor(0xFF,0xFF,0xFF) if i == 0 else INK
            if i == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = PRIMARY
    # Big gap callout
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(8.5), Inches(1.4), Inches(4.5), Inches(4.7))
    card.fill.solid(); card.fill.fore_color.rgb = ACCENT; card.line.fill.background()
    white = RGBColor(0xFF,0xFF,0xFF)
    add_textbox(s, Inches(8.5), Inches(1.6), Inches(4.5), Inches(0.6),
                "Dormant", size=20, bold=True, color=white, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(8.5), Inches(2.1), Inches(4.5), Inches(1.0),
                "90%", size=56, bold=True, color=white, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(8.5), Inches(3.2), Inches(4.5), Inches(0.6),
                "Active loyal", size=20, bold=True, color=white, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(8.5), Inches(3.7), Inches(4.5), Inches(1.0),
                "23%", size=56, bold=True, color=white, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(8.5), Inches(5.0), Inches(4.5), Inches(0.8),
                "67-point churn gap\n→ segment-aware retention works",
                size=14, color=white, align=PP_ALIGN.CENTER)
    add_footer(s, 17)
    add_speaker_notes(s, "This is THE slide. Pause here. Explain the gap is the synthesis insight.")

    # ---------- 18. Cross — per-cluster rules + heatmap ----------------- #
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Loyal Segment Buys in Colour Pairs",
                  "Per-segment FP-Growth recovers cross-sell opportunities")
    # Heatmap is 1.6:1 — constrain by HEIGHT so it doesn't run into the recommendation
    add_image(s, FIGURES_DIR / "06_segment_products_heatmap.png",
              Inches(0.4), Inches(1.3), height=Inches(4.5))
    add_bullets(s, Inches(8.0), Inches(1.6), Inches(5.2), Inches(4),
                ["Pink Lunchbag → Red\nLunchbag (lift 7.4)",
                 "Pink Jumbo Bag → Red\nJumbo Bag (lift 6.4)",
                 "Red Heart T-light →\nWhite Heart T-light (5.3)"], size=14)
    add_textbox(s, Inches(0.5), Inches(6.05), Inches(12.5), Inches(0.95),
                "Recommendation: target the 23% churners in the loyal segment "
                "with colour-pair bundles. Use win-back campaigns (not bundles) "
                "for the dormant majority.",
                size=14, bold=True, color=PRIMARY)
    add_footer(s, 18)

    # ---------- 19. Limitations ----------------------------------------- #
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Limitations & Future Work",
                  "Honest caveats — and where we'd go next")
    add_textbox(s, Inches(0.5), Inches(1.3), Inches(6), Inches(0.4),
                "Limitations", size=18, bold=True, color=PRIMARY)
    add_bullets(s, Inches(0.5), Inches(1.8), Inches(6), Inches(4),
                [f"Leakage fix: supervised features built strictly before {tokens['churn_cutoff']} — "
                 f"best AUC now a realistic {tokens['classification_best_auc']} (full-history features were inflated)",
                 f"Recency dominant but not sufficient — ablation without it drops AUC to {tokens['without_recency_auc']}; "
                 f"baselines: majority {tokens['majority_auc']}, recency-only {tokens['recency_auc']}",
                 "Single retailer, UK-centric, 2 yrs — generalisability untested",
                 "No demographic attributes available",
                 "Dormant segment too sparse for in-cluster rules",
                 "Anomaly threshold purely percentile-based"], size=13)
    add_textbox(s, Inches(6.8), Inches(1.3), Inches(6), Inches(0.4),
                "Future Work", size=18, bold=True, color=PRIMARY)
    add_bullets(s, Inches(6.8), Inches(1.8), Inches(6), Inches(4),
                ["Transformer on purchase sequences (next-basket + churn jointly)",
                 "Uplift modelling: who responds best to retention offers",
                 "Cross-retailer transfer to test segment-rule stability",
                 "Domain-labelled fraud / VIP ground truth for anomalies"], size=13)
    add_footer(s, 19)

    # ---------- 20. Conclusion ------------------------------------------- #
    s = prs.slides.add_slide(blank)
    add_title_bar(s, "Conclusion & Thank You", "Q&A — happy to discuss any slide")
    add_textbox(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.5),
                "One pipeline, four course topics, one business insight.",
                size=22, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
    add_bullets(s, Inches(0.6), Inches(2.2), Inches(12), Inches(4),
                ["Segmentation: K-means k=2, Silhouette 0.42",
                 f"Classification: {best['Model']} AUC {best['AUC']:.3f} · MLP AUC {mlp['AUC']:.3f} (leakage-safe pre-cutoff features)",
                 "Association: 17 rules, FP-Growth 1.5× faster, top lift 26.8",
                 "★ Synthesis: 67-pt churn gap → colour-pair retention bundles"],
                size=18)
    add_textbox(s, Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.6),
                "Code, report & slides — Team 4 GitHub repository",
                size=14, color=SUBINK, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.6),
                "Team 4 — Cuong · Danh · Dat · Dinh    •    UIT, VNU-HCM    •    02 June 2026",
                size=13, color=SUBINK, align=PP_ALIGN.CENTER)
    add_footer(s, 20)

    return prs


def main():
    deck = build_deck()
    out = REPORTS_DIR / "slides.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    deck.save(out)
    logger.info("Wrote %s  (%.1f KB, %d slides)", out, out.stat().st_size / 1024, len(deck.slides))


if __name__ == "__main__":
    main()
