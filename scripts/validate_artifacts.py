"""Submission-grade artifact validator.

Confirms the pipeline's outputs are present, schema-correct, mutually
consistent, leakage-free, and free of stale metrics / placeholders. Exits 0
with a concise summary when everything checks out, non-zero otherwise — safe to
wire into CI or a pre-submission check.

    python scripts/validate_artifacts.py
"""
from __future__ import annotations

import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))

import pandas as pd  # noqa: E402

from src import artifacts as A  # noqa: E402

PLACEHOLDERS = ["<repo-url>", "<your-org>", "TODO_GITHUB_URL", "TODO_GDRIVE_LINK"]
# Generated data artifacts: scanned for placeholders only (NOT for "0.99x" stale
# numbers, which would false-positive on legitimate support/confidence values).
ARTIFACT_TEXT_GLOBS = ["reports/*.csv", "reports/manifest.json"]
# Submission-facing narrative/source files: scanned for BOTH placeholders and the
# specific inflated pre-leakage-fix metrics. `Claude/` is intentionally excluded
# (the enhancement notes legitimately discuss the old numbers).
SUBMISSION_TEXT_GLOBS = [
    "README.md", "README.template.md",
    "overleaf/*.tex", "overleaf/*.md",
    "src/build_slides.py",
]
STALE_METRICS = ["0.998", "0.9981", "0.990", "0.9900"]

REQUIRED = [
    "transactions_clean", "rfm_features", "churn_labels", "classification_features",
    "classical_results", "baseline_results", "ablation_results", "calibration_results",
    "cluster_stability", "association_rules", "mlp_vs_classical", "manifest",
    "customer_churn_scores", "customer_decision_table",
    "customer_decision_summary", "top_priority_customers",
]
SCHEMA_KEYS = [
    "rfm_features", "churn_labels", "classification_features", "classical_results",
    "baseline_results", "ablation_results", "calibration_results", "association_rules",
    "customer_churn_scores", "customer_decision_table",
    "customer_decision_summary", "top_priority_customers",
]


class Validator:
    def __init__(self) -> None:
        self.errors: list[str] = []
        self.passed: list[str] = []

    def check(self, ok: bool, label: str, detail: str = "") -> None:
        if ok:
            self.passed.append(label)
        else:
            self.errors.append(f"{label}: {detail}" if detail else label)

    # --- individual checks -------------------------------------------------- #
    def files_exist(self) -> None:
        missing = A.missing_artifacts(REQUIRED)
        self.check(not missing, "required files exist", f"missing {missing}")

    def schemas(self) -> None:
        for key in SCHEMA_KEYS:
            try:
                miss = A.check_csv_schema(key)
                self.check(not miss, f"schema[{key}]", f"missing cols {miss}")
            except FileNotFoundError:
                self.check(False, f"schema[{key}]", "file missing")

    def consistency(self) -> None:
        try:
            feats = pd.read_csv(A.ARTIFACTS["classification_features"])
            churn = pd.read_csv(A.ARTIFACTS["churn_labels"])
            same = set(feats["CustomerID"]) == set(churn["CustomerID"])
            self.check(same, "classification_features == churn_labels customer set",
                       "customer sets differ")
        except Exception as e:  # noqa: BLE001
            self.check(False, "feature/label consistency", str(e))

    def no_missing_values(self) -> None:
        for key in ("classification_features", "churn_labels", "classical_results"):
            try:
                df = pd.read_csv(A.ARTIFACTS[key])
                self.check(not df.isnull().any().any(),
                           f"no missing values[{key}]", "found NaNs")
            except Exception as e:  # noqa: BLE001
                self.check(False, f"no missing values[{key}]", str(e))

    def leakage_check(self) -> None:
        """Recency built pre-cutoff must be non-negative, and AUC realistic."""
        try:
            feats = pd.read_csv(A.ARTIFACTS["classification_features"])
            self.check((feats["Recency"] >= 0).all(),
                       "leakage: Recency >= 0 (pre-cutoff snapshot)", "negative recency")
        except Exception as e:  # noqa: BLE001
            self.check(False, "leakage: recency", str(e))
        try:
            res = pd.read_csv(A.ARTIFACTS["classical_results"])
            best_auc = float(res["AUC"].max())
            self.check(0.6 <= best_auc <= 0.95,
                       "leakage: best AUC in realistic range [0.6, 0.95]",
                       f"AUC={best_auc} suggests leakage or a bug")
        except Exception as e:  # noqa: BLE001
            self.check(False, "leakage: AUC range", str(e))

    def no_placeholders_in_artifacts(self) -> None:
        # Placeholders must never appear in machine-generated data artifacts.
        offenders = []
        for pattern in ARTIFACT_TEXT_GLOBS:
            for path in ROOT.glob(pattern):
                text = path.read_text(encoding="utf-8", errors="ignore")
                for ph in PLACEHOLDERS:
                    if ph in text:
                        offenders.append(f"{path.name}:{ph}")
        self.check(not offenders, "no placeholders in data artifacts",
                   ", ".join(offenders))

    def no_stale_metrics(self) -> None:
        # Submission-facing files must not quote the inflated pre-fix metrics or
        # leftover placeholders. (Generated CSVs are excluded to avoid false
        # positives on legitimate support/confidence values.)
        offenders = []
        for pattern in SUBMISSION_TEXT_GLOBS:
            for path in ROOT.glob(pattern):
                text = path.read_text(encoding="utf-8", errors="ignore")
                for token in STALE_METRICS + PLACEHOLDERS:
                    if token in text:
                        offenders.append(f"{path.relative_to(ROOT)}:{token}")
        self.check(not offenders, "no stale leaky metrics / placeholders in submission files",
                   ", ".join(offenders))

    def rendered_docs_match_snapshot(self) -> None:
        try:
            from src.submission_snapshot import load_submission_snapshot
            s = load_submission_snapshot(ROOT)
        except Exception as e:  # noqa: BLE001
            self.check(False, "snapshot loads for doc check", str(e))
            return
        expected = [s.best_model, s.best_auc_3dp, s.mlp_auc_3dp, s.churn_cutoff]
        for rel in ("README.md", "overleaf/main.tex"):
            path = ROOT / rel
            if not path.exists():
                self.check(False, f"{rel} exists", "missing")
                continue
            text = path.read_text(encoding="utf-8", errors="ignore")
            missing = [v for v in expected if v not in text]
            self.check(not missing, f"{rel} reflects current metrics", f"missing {missing}")

    def pdf_freshness(self) -> None:
        # PDFs are optional here (no LaTeX/PowerPoint toolchain assumed); but if a
        # slides.pdf is present it must not be older than the slides.pptx it derives
        # from, or it would ship stale numbers.
        pptx = A.REPORTS_DIR / "slides.pptx"
        pdf = A.REPORTS_DIR / "slides.pdf"
        if pdf.exists() and pptx.exists():
            self.check(pdf.stat().st_mtime >= pptx.stat().st_mtime,
                       "slides.pdf is at least as fresh as slides.pptx",
                       "re-export slides.pdf from the updated slides.pptx")
        else:
            self.check(True, "slides.pdf freshness (optional — pdf not present)", "")

    def manifest_sane(self) -> None:
        try:
            m = A.read_manifest()
            self.check(m.get("seed") == 42, "manifest seed == 42", str(m.get("seed")))
            self.check(m.get("churn_cutoff") == "2011-09-01",
                       "manifest churn_cutoff == 2011-09-01", str(m.get("churn_cutoff")))
        except Exception as e:  # noqa: BLE001
            self.check(False, "manifest readable", str(e))

    def decisions(self) -> None:
        try:
            table = pd.read_csv(A.ARTIFACTS["customer_decision_table"])
            feats = pd.read_csv(A.ARTIFACTS["classification_features"])
        except Exception as e:  # noqa: BLE001
            self.check(False, "decision table readable", str(e))
            return
        self.check(len(table) == len(feats),
                   "decision rows == scored customers", f"{len(table)} vs {len(feats)}")
        self.check(table["recommended_action"].notna().all(),
                   "no null recommended_action", "found nulls")
        self.check(table["priority_score"].notna().all(),
                   "no null priority_score", "found nulls")
        self.check(sorted(table["priority_rank"]) == list(range(1, len(table) + 1)),
                   "priority_rank is a contiguous 1..N", "ranks not contiguous")

        m = A.read_manifest()
        self.check("decision_table_customers" in m,
                   "manifest has decision metadata", "missing decision_* fields")
        self.check(m.get("decision_table_customers") == len(table),
                   "manifest decision_table_customers matches table",
                   f"{m.get('decision_table_customers')} vs {len(table)}")
        src = set(table["source_run_generated_at"].astype(str))
        self.check(len(src) == 1 and str(m.get("decision_source_run_generated_at")) in src,
                   "source_run_generated_at matches manifest", f"table={src}")

        try:
            top = pd.read_csv(A.ARTIFACTS["top_priority_customers"])
            self.check(set(top["CustomerID"]).issubset(set(table["CustomerID"])),
                       "top-priority customers subset of table", "not a subset")
        except Exception as e:  # noqa: BLE001
            self.check(False, "top-priority readable", str(e))

    def slides_present(self) -> None:
        pptx = A.REPORTS_DIR / "slides.pptx"
        if not pptx.exists():
            self.check(True, "slides.pptx (optional — skipped)", "")
            return
        try:
            from pptx import Presentation
            n = len(Presentation(str(pptx)).slides)
            self.check(n == 20, "slide deck has 20 slides", f"found {n}")
        except Exception as e:  # noqa: BLE001
            self.check(False, "slide deck readable", str(e))

    # --- driver ------------------------------------------------------------- #
    def run(self) -> int:
        for fn in (self.files_exist, self.schemas, self.consistency,
                   self.no_missing_values, self.leakage_check,
                   self.no_placeholders_in_artifacts, self.no_stale_metrics,
                   self.rendered_docs_match_snapshot, self.manifest_sane,
                   self.decisions,
                   self.slides_present, self.pdf_freshness):
            fn()

        print(f"\n{'=' * 60}")
        print(f"  ARTIFACT VALIDATION  —  {len(self.passed)} passed, {len(self.errors)} failed")
        print(f"{'=' * 60}")
        for p in self.passed:
            print(f"  PASS  {p}")
        for e in self.errors:
            print(f"  FAIL  {e}")
        print(f"{'=' * 60}")
        if self.errors:
            print("VALIDATION FAILED")
            return 1
        print("ALL CHECKS PASSED")
        return 0


if __name__ == "__main__":
    raise SystemExit(Validator().run())
